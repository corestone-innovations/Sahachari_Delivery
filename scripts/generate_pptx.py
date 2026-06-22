from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os

root = os.path.dirname(os.path.dirname(__file__))
assets = os.path.join(root, 'assets', 'images')
logo = os.path.join(assets, 'logo.png')
splash = os.path.join(assets, 'splashscreen.png')

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle, logo_path=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    subtitle_tf = slide.placeholders[1]
    title_tf.text = title
    subtitle_tf.text = subtitle
    if logo_path and os.path.exists(logo_path):
        left = Inches(10.8)
        top = Inches(0.3)
        height = Inches(1.0)
        slide.shapes.add_picture(logo_path, left, top, height=height)
    return slide


def add_bullet_slide(title, bullets, notes=None, image_path=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    if image_path and os.path.exists(image_path):
        left = Inches(9.6)
        top = Inches(1.2)
        height = Inches(3.6)
        slide.shapes.add_picture(image_path, left, top, height=height)
    return slide

slides = [
    {
        'title': 'Sahachari Delivery',
        'subtitle': 'Last-mile delivery app — Available Jobs, Orders & Payments',
        'logo': logo
    },
    {
        'title': 'Introduction',
        'bullets': [
            'Mobile-first delivery platform for riders, merchants and customers',
            'Supports browsing available delivery jobs, accepting work, and managing orders',
            'Built using Expo / React Native with backend REST API integration'
        ],
        'notes': 'Introduce the app and the main user groups it supports.'
    },
    {
        'title': 'Objectives of the Project',
        'bullets': [
            'Deliver a rider app for discovering READY deliveries quickly',
            'Allow riders to accept jobs and update order status in one place',
            'Capture payment transactions directly from the rider workflow',
            'Maintain up-to-date rider profile and contact data'
        ],
        'notes': 'Emphasize the practical goals already implemented in the code.'
    },
    {
        'title': 'Presentation Agenda',
        'bullets': [
            'Business context and app goals',
            'Implemented technical architecture',
            'Core screens and features',
            'Reliability and current limitations'
        ]
    },
    {
        'title': 'Business Understanding',
        'bullets': [
            'Target users: delivery riders, merchants listing orders, customers receiving packages',
            'Need for faster local order fulfillment and simple rider onboarding',
            'Value: better visibility into jobs, payments, and rider profile updates'
        ]
    },
    {
        'title': 'The Problem',
        'bullets': [
            'No central app for riders to view ready-to-deliver orders',
            'Manual delivery coordination slows down the process',
            'Riders need a reliable way to confirm order status and payment details'
        ]
    },
    {
        'title': 'Why This Project?',
        'bullets': [
            'Simplifies rider operations with a structured job list',
            'Improves delivery throughput by guiding status updates',
            'Supports offline-screen style user flow with clear success/failure alerts'
        ]
    },
    {
        'title': 'Development Methodology',
        'bullets': [
            'Iterative build: feature-first screens and API integrations',
            'Use React Query to cache and refresh API data smoothly',
            'Implement user feedback through alerts and error handling'
        ]
    },
    {
        'title': 'Technologies Used',
        'bullets': [
            'React Native + Expo for cross-platform mobile UI',
            'TypeScript for typed components and API models',
            '@tanstack/react-query for server state and caching',
            'expo-linear-gradient and icons for polished visual design'
        ]
    },
    {
        'title': 'Architecture & Data Flow',
        'bullets': [
            'Authentication token managed by AuthContext',
            'Central apiRequest wrapper handles REST calls and headers',
            'Available jobs, orders, and profile data fetched with React Query',
            'Mutations invalidate relevant queries to refresh UI data'
        ]
    },
    {
        'title': 'Core System Modules (Implemented)',
        'bullets': [
            'Available Jobs: fetch `/delivery/orders?status=READY` and display job cards',
            'Accept Job: POST `/delivery/orders/{id}/accept` with mutation feedback',
            'My Orders: load `/delivery/orders?mine=true` and show order status flow',
            'Profile: update user contact, address, photo, and serviceable pincodes'
        ]
    },
    {
        'title': 'Available Jobs Screen',
        'bullets': [
            'Displays ready orders with zone badge, store name, and full address details',
            'Includes delivery amount, status badge, and accept button per job',
            'Supports pull-to-refresh and loading state UI',
            'Accept action invalidates job lists and alerts rider on success/failure'
        ],
        'image': splash
    },
    {
        'title': 'My Orders & Payment Flow',
        'bullets': [
            'Fetches rider-specific jobs and shows current status steps',
            'Status actions available: pickup, deliver, fail',
            'Payment flow fetches UPI ID from `/upi-collection/checkout/{checkoutId}` or `/upi-collection/order/{orderId}`',
            'Creates payment transaction via `/payment-transactions` and displays QR payment details'
        ]
    },
    {
        'title': 'Profile & Rider Management',
        'bullets': [
            'Loads rider profile from `/users/me`',
            'Editable fields: mobile number, address, and profile picture',
            'Shows serviceable pincodes and supports logout',
            'Validates input and provides success/error alerts for updates'
        ]
    },
    {
        'title': 'Error Handling & Reliability',
        'bullets': [
            'Fallback fetch for available jobs when `/delivery/orders?status=READY` fails',
            'Clear Alert dialogs for accept, payment, and profile operations',
            'React Query cache invalidation keeps displayed data current',
            'Pending states disable buttons and show loading indicators'
        ]
    },
    {
        'title': 'Current App Limitations',
        'bullets': [
            'No dedicated financial dashboard screen in the current app',
            'Cash & bank reconciliation module not implemented yet',
            'Registry/family management is outside current delivery scope'
        ]
    },
    {
        'title': 'Questions?',
        'bullets': [
            'Would you like a demo of Available Jobs, My Orders, or Profile screen?',
            'Which part of the app would you like to discuss in more detail?'
        ]
    },
    {
        'title': 'Image Sources',
        'bullets': [
            'Project asset images from the app repository',
            'Add live emulator screenshots for the final teacher presentation'
        ]
    }
]

# Build slides
add_title_slide(slides[0]['title'], slides[0]['subtitle'], slides[0]['logo'])
for s in slides[1:]:
    add_bullet_slide(s['title'], s.get('bullets', []), notes=s.get('notes'), image_path=s.get('image'))

out_path = os.path.join(root, 'sahachari_delivery_presentation_v2.pptx')
prs.save(out_path)
print('Created:', out_path)
